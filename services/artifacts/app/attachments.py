"""Attachments-as-context: extract text from uploaded files, chunk it, store it,
and serve relevant passages so the agent can use file content when generating.

Phase 1: extract -> store (SQLite) -> return whole (capped) text on retrieve.
Phase 2 (gated, not implemented here): embeddings + RAG for large docs. The
`chunks.embedding` column exists but stays NULL in Phase 1.
"""
import csv
import io
import os
import sqlite3
import threading
import uuid

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

MAX_UPLOAD_BYTES = 10 * 1024 * 1024
# Docs at or below this length are stored whole (one chunk); larger docs are chunked.
WHOLE_CAP = 8000

# Module-level SQLite path from env with a sensible default alongside the service.
_DB = os.environ.get(
    "ATTACHMENTS_DB",
    os.path.join(os.path.dirname(__file__), "..", "attachments.db"),
)
_lock = threading.Lock()


def _conn() -> sqlite3.Connection:
    """Open a connection and ensure tables exist (created on demand)."""
    con = sqlite3.connect(_DB)
    con.execute(
        "CREATE TABLE IF NOT EXISTS attachments("
        "doc_id TEXT PRIMARY KEY, name TEXT, mime TEXT, chars INT, chunk_count INT, created_at REAL)"
    )
    con.execute(
        "CREATE TABLE IF NOT EXISTS chunks(doc_id TEXT, ord INT, text TEXT, embedding BLOB)"
    )
    return con


def extract_text(data: bytes, mime: str, name: str) -> str:
    n = name.lower()
    if n.endswith((".md", ".txt")) or mime.startswith("text/"):
        return data.decode("utf-8", "replace")
    if n.endswith(".csv"):
        rows = list(csv.reader(io.StringIO(data.decode("utf-8", "replace"))))
        return "\n".join(", ".join(r) for r in rows)
    if n.endswith(".docx"):
        return "\n".join(p.text for p in Document(io.BytesIO(data)).paragraphs if p.text)
    if n.endswith(".xlsx"):
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        out = []
        for ws in wb.worksheets:
            out.append(f"# {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    out.append(", ".join(cells))
        return "\n".join(out)
    if n.endswith(".pptx"):
        out = []
        for i, s in enumerate(Presentation(io.BytesIO(data)).slides, 1):
            out.append(f"# Slide {i}")
            for sh in s.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip():
                    out.append(sh.text_frame.text)
        return "\n".join(out)
    if n.endswith(".pdf"):
        return "\n".join((pg.extract_text() or "") for pg in PdfReader(io.BytesIO(data)).pages)
    raise ValueError(f"unsupported file type: {name}")


def chunk_text(text: str, size: int = 1000, overlap: int = 150) -> list[str]:
    if len(text) <= size:
        return [text]
    step = max(1, size - overlap)
    return [text[i:i + size] for i in range(0, len(text), step)]


def store_attachment(data: bytes, mime: str, name: str, now: float) -> dict:
    if len(data) > MAX_UPLOAD_BYTES:
        raise ValueError("file too large")
    text = extract_text(data, mime, name)
    chunks = [text] if len(text) <= WHOLE_CAP else chunk_text(text)
    doc_id = "doc-" + uuid.uuid4().hex[:12]
    with _lock, _conn() as con:
        con.execute(
            "INSERT INTO attachments VALUES (?,?,?,?,?,?)",
            (doc_id, name, mime, len(text), len(chunks), now),
        )
        con.executemany(
            "INSERT INTO chunks(doc_id, ord, text, embedding) VALUES (?,?,?,NULL)",
            [(doc_id, i, c) for i, c in enumerate(chunks)],
        )
    return {
        "doc_id": doc_id,
        "name": name,
        "chars": len(text),
        "chunk_count": len(chunks),
        "preview": text[:200],
    }


# Phase 1: return each doc's whole text (chunks concatenated in order), capped so
# the combined context stays within a sane budget. Embeddings are ignored (NULL).
# Phase 2 (gated) would embed the query and cosine-rank chunks for large docs.
TOTAL_CAP = 24000


def retrieve_context(doc_ids: list[str], query: str, k: int = 6) -> list[dict]:
    passages, budget = [], TOTAL_CAP
    with _conn() as con:
        for doc_id in doc_ids:
            meta = con.execute(
                "SELECT name FROM attachments WHERE doc_id=?", (doc_id,)
            ).fetchone()
            if not meta:
                continue
            rows = con.execute(
                "SELECT text FROM chunks WHERE doc_id=? ORDER BY ord", (doc_id,)
            ).fetchall()
            text = "\n".join(r[0] for r in rows)[:budget]
            if not text:
                continue
            passages.append({"doc_id": doc_id, "name": meta[0], "text": text})
            budget -= len(text)
            if budget <= 0:
                break
    return passages
