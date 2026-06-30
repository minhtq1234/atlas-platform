"""DeckContent -> .pptx bytes (python-pptx)."""
import io

from pptx import Presentation
from pptx.util import Inches, Pt

from . import brand
from .models import DeckContent, Slide


def _bg(slide, hex6: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = brand.pptx_rgb(hex6)


def _text(slide, left, top, width, height, text, *, size, color, bold=False,
          font=brand.UI_FONT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = brand.pptx_rgb(color)
    return box


def _bullets(slide, left, top, width, height, bullets, *, color):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = "•  " + b
        run.font.size = Pt(18)
        run.font.name = brand.UI_FONT
        run.font.color.rgb = brand.pptx_rgb(color)


def _cover(slide, deck: DeckContent, s: Slide):
    _bg(slide, brand.INK)
    _text(slide, 0.9, 0.8, 11, 0.5, deck.eyebrow.upper(), size=12, color=brand.CORAL, bold=True)
    _text(slide, 0.9, 2.6, 11.5, 2.5, s.title, size=44, color=brand.WHITE, bold=True, font=brand.SERIF_FONT)
    _text(slide, 0.9, 5.0, 11, 0.6, s.subtitle or deck.subtitle, size=16, color="C7CEE8")


def _content(slide, deck: DeckContent, s: Slide):
    _bg(slide, brand.WHITE)
    _text(slide, 0.9, 0.7, 11, 0.4, deck.eyebrow.upper(), size=11, color=brand.INDIGO, bold=True)
    _text(slide, 0.9, 1.2, 11.5, 1.0, s.title, size=30, color=brand.INK, bold=True)
    if s.bullets:
        _bullets(slide, 1.0, 2.6, 11, 4.0, s.bullets, color=brand.INK)


def build_deck(content: DeckContent, name: str) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = name
    blank = prs.slide_layouts[6]

    for s in content.slides:
        slide = prs.slides.add_slide(blank)
        if s.isCover:
            _cover(slide, content, s)
        else:
            _content(slide, content, s)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
