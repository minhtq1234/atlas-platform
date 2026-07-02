import io

from docx import Document
from fastapi.testclient import TestClient
from openpyxl import load_workbook
from pptx import Presentation

from app.docx_builder import build_doc
from app.main import app
from app.models import DeckContent, DocContent, SheetContent
from app.pptx_builder import build_deck
from app.xlsx_builder import build_sheet

OOXML_MAGIC = b"PK\x03\x04"


def _formula_cells(ws):
    return [c.value for row in ws.iter_rows() for c in row
            if isinstance(c.value, str) and c.value.startswith("=")]


def test_doc_is_valid_docx_with_title_and_paragraphs():
    c = DocContent(
        kind="Doc", eyebrow="MEMO · HR", title="Q2 Headcount Review",
        meta="30 June 2026", paragraphs=["First para.", "Second para."],
        bars=[{"label": "Apr", "value": 0.5}], callout={"value": "+18", "label": "NET"},
    )
    data = build_doc(c, "Headcount Memo")
    assert data[:4] == OOXML_MAGIC
    doc = Document(io.BytesIO(data))
    texts = "\n".join(p.text for p in doc.paragraphs)
    assert "Q2 Headcount Review" in texts
    assert "First para." in texts


def test_sheet_is_valid_xlsx_with_live_sum_formula():
    c = SheetContent(
        kind="Sheet", title="Headcount Model", columns=["Org unit", "Net", "EoP"],
        rows=[["TSE", "+18", 342], ["Platform", "+11", 208], ["Total", "+29", 550]],
    )
    data = build_sheet(c, "Headcount Model")
    assert data[:4] == OOXML_MAGIC
    wb = load_workbook(io.BytesIO(data))
    ws = wb.active
    assert ws["A2"].value == "Org unit"
    # a real SUM formula must exist somewhere in the last column
    formulas = [cell.value for row in ws.iter_rows() for cell in row
                if isinstance(cell.value, str) and cell.value.startswith("=SUM")]
    assert formulas, "expected a live =SUM formula"


def test_deck_is_valid_pptx_with_one_slide_per_input():
    c = DeckContent(
        kind="Deck", eyebrow="ATLAS · BOARD DECK", title="Q2 People Review",
        subtitle="Headcount & hiring",
        slides=[
            {"isCover": True, "title": "Q2 People Review", "subtitle": "Headcount & hiring"},
            {"title": "Headcount", "bullets": ["1,248 total", "+18 net"]},
        ],
    )
    data = build_deck(c, "People Review")
    assert data[:4] == OOXML_MAGIC
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 2
    all_text = " ".join(
        run.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        for run in para.runs
    )
    assert "Q2 People Review" in all_text
    assert "1,248 total" in all_text


def test_deck_exports_speaker_notes_and_section_layout():
    c = DeckContent(
        kind="Deck", eyebrow="ATLAS · BOARD", title="Q2", subtitle="x",
        slides=[
            {"isCover": True, "title": "Q2"},
            {"title": "Part Two", "layout": "section"},
            {"title": "Growth up 40%", "bullets": ["Enterprise led"], "notes": "Mention the Q2 deal."},
        ],
    )
    data = build_deck(c, "Q2 Deck")
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 3
    # speaker notes land on the third slide's notes page
    assert "Q2 deal" in prs.slides[2].notes_slide.notes_text_frame.text


def test_sheet_neutralizes_formula_injection():
    # A string cell starting with '=' must NOT become a live formula.
    c = SheetContent(
        kind="Sheet", title="x", columns=["Item", "Val"],
        rows=[["=HYPERLINK(\"http://evil\",\"x\")", 5], ["safe", 7]],
    )
    wb = load_workbook(io.BytesIO(build_sheet(c, "s")))
    formulas = _formula_cells(wb.active)
    # the only formula present is our own trusted SUM
    assert all(f.startswith("=SUM") for f in formulas), formulas


def test_sheet_sum_correct_when_total_row_not_last():
    # Total in the MIDDLE must not truncate/invert the SUM range.
    c = SheetContent(
        kind="Sheet", title="x", columns=["Unit", "N"],
        rows=[["TSE", 10], ["Total", 99], ["Platform", 20]],
    )
    wb = load_workbook(io.BytesIO(build_sheet(c, "s")))
    sums = [f for f in _formula_cells(wb.active) if f.startswith("=SUM")]
    assert sums, "expected a SUM"
    # data rows are 3 (TSE) and 5 (Platform); Total at row 4 excluded
    assert sums[0] == "=SUM(B3,B5)", sums[0]


def test_zero_column_sheet_returns_422_not_500():
    client = TestClient(app)
    r = client.post("/export", json={"name": "x", "content": {
        "kind": "Sheet", "title": "t", "columns": [], "rows": []}})
    assert r.status_code == 422, r.status_code


def test_vietnamese_filename_survives_in_content_disposition():
    client = TestClient(app)
    r = client.post("/export", json={"name": "Báo cáo Tài chính", "content": {
        "kind": "Doc", "eyebrow": "E", "title": "T", "meta": "m", "paragraphs": ["p"]}})
    assert r.status_code == 200
    cd = r.headers["content-disposition"]
    assert "filename*=UTF-8''" in cd
    from urllib.parse import quote
    assert quote("Báo cáo Tài chính.docx", safe="") in cd


def test_docx_renders_sections():
    c = DocContent.model_validate({
        "kind": "Doc", "eyebrow": "E", "title": "T", "meta": "m",
        "sections": [{"heading": "Requirements", "blocks": [
            {"type": "paragraph", "text": "intro"},
            {"type": "bullets", "items": ["a", "b"]},
            {"type": "table", "columns": ["ID", "Req"], "rows": [["FR-1", "Login"]]},
        ]}],
    })
    data = build_doc(c, "Doc")
    doc = Document(io.BytesIO(data))
    text = "\n".join(p.text for p in doc.paragraphs)
    assert "Requirements" in text and "intro" in text
    assert len(doc.tables) == 1 and doc.tables[0].rows[0].cells[0].text == "ID"
